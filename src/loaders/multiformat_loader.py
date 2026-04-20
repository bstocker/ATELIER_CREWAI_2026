from __future__ import annotations

from pathlib import Path
from typing import Iterable
import csv
import xml.etree.ElementTree as ET

from pypdf import PdfReader
from PIL import Image
import pytesseract
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from src.loaders.admission_ocr_loader import extract_text_with_simple_ocr


# --------------------------------------------------
# EXTENSIONS
# --------------------------------------------------

TEXT_EXTENSIONS = {".md", ".txt", ".csv", ".xml"}

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".xlsx",
    ".xls",
    ".docx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
    ".xml",
    ".csv",
    ".md",
    ".txt",
    ".pptx",
    ".wav",
    ".mp3",
    ".m4a",
}


# --------------------------------------------------
# OUTILS GÉNÉRAUX
# --------------------------------------------------

def _is_text_usable(text: str, min_chars: int = 80, min_words: int = 15) -> bool:
    if not text:
        return False

    text = text.strip()
    if len(text) < min_chars:
        return False

    if len(text.split()) < min_words:
        return False

    return True


# --------------------------------------------------
# TXT / MD
# --------------------------------------------------

def _load_text(file_path: Path) -> str:
    return file_path.read_text(encoding="utf-8", errors="ignore")


# --------------------------------------------------
# CSV
# --------------------------------------------------

def _load_csv(file_path: Path) -> str:
    rows = []

    with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))

    return "\n".join(rows)


# --------------------------------------------------
# XML
# --------------------------------------------------

def _load_xml(file_path: Path) -> str:
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        return ET.tostring(root, encoding="unicode")
    except Exception as exc:
        return f"[ERREUR XML : {exc}]"


# --------------------------------------------------
# PDF STRUCTURÉ
# --------------------------------------------------

def _extract_text_from_pdf(file_path: Path) -> str:
    reader = PdfReader(str(file_path))
    pages = []

    for i, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
            pages.append(f"\n--- PAGE {i} ---\n{text.strip()}")
        except Exception as exc:
            pages.append(f"\n--- PAGE {i} ---\n[ERREUR PDF : {exc}]")

    return "\n".join(pages).strip()


# --------------------------------------------------
# IMAGE OCR
# --------------------------------------------------

def _load_image_with_ocr(file_path: Path, lang: str = "fra") -> str:
    try:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img, lang=lang)
    except Exception as exc:
        return f"[ERREUR OCR IMAGE : {exc}]"


# --------------------------------------------------
# DOCX
# --------------------------------------------------

def _load_docx(file_path: Path) -> str:
    try:
        doc = Document(str(file_path))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
    except Exception as exc:
        return f"[ERREUR DOCX : {exc}]"


# --------------------------------------------------
# XLSX / XLS
# --------------------------------------------------

def _load_excel(file_path: Path) -> str:
    try:
        wb = load_workbook(filename=file_path, data_only=True)
        output = []

        for ws in wb.worksheets:
            output.append(f"\n### FEUILLE : {ws.title}")

            for row in ws.iter_rows(values_only=True):
                values = [str(v) if v is not None else "" for v in row]
                if any(values):
                    output.append(" | ".join(values))

        return "\n".join(output)

    except Exception as exc:
        return f"[ERREUR EXCEL : {exc}]"


# --------------------------------------------------
# PPTX
# --------------------------------------------------

def _load_pptx(file_path: Path) -> str:
    try:
        prs = Presentation(str(file_path))
        slides = []

        for i, slide in enumerate(prs.slides, start=1):
            slides.append(f"\n### SLIDE {i}")

            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    txt = shape.text.strip()
                    if txt:
                        slides.append(txt)

        return "\n".join(slides)

    except Exception as exc:
        return f"[ERREUR PPTX : {exc}]"


# --------------------------------------------------
# AUDIO (placeholder)
# --------------------------------------------------

def _load_audio(file_path: Path) -> str:
    return f"[FICHIER AUDIO DÉTECTÉ : {file_path.name} - transcription non activée]"


# --------------------------------------------------
# CHARGEMENT UNIQUE
# --------------------------------------------------

def load_single_file(
    file_path: str | Path,
    enable_ocr: bool = False,
    ocr_lang: str = "fra",
) -> str:

    file_path = Path(file_path)
    suffix = file_path.suffix.lower()

    if suffix not in SUPPORTED_EXTENSIONS:
        raise ValueError(f"Format non supporté : {suffix}")

    # TXT / MD
    if suffix in {".txt", ".md"}:
        return _load_text(file_path)

    # CSV
    if suffix == ".csv":
        return _load_csv(file_path)

    # XML
    if suffix == ".xml":
        return _load_xml(file_path)

    # PDF
    if suffix == ".pdf":
        native_text = _extract_text_from_pdf(file_path)

        if _is_text_usable(native_text):
            return native_text

        if enable_ocr:
            return extract_text_with_simple_ocr(
                file_path=file_path,
                lang=ocr_lang
            )

        return native_text

    # Images
    if suffix in {".png", ".jpg", ".jpeg", ".webp"}:
        if enable_ocr:
            return _load_image_with_ocr(file_path, lang=ocr_lang)
        return f"[IMAGE DÉTECTÉE : {file_path.name}]"

    # DOCX
    if suffix == ".docx":
        return _load_docx(file_path)

    # Excel
    if suffix in {".xlsx", ".xls"}:
        return _load_excel(file_path)

    # PPTX
    if suffix == ".pptx":
        return _load_pptx(file_path)

    # Audio
    if suffix in {".wav", ".mp3", ".m4a"}:
        return _load_audio(file_path)

    return ""


# --------------------------------------------------
# PARCOURS RÉPERTOIRE
# --------------------------------------------------

def iter_files(directory: str | Path) -> Iterable[Path]:
    directory = Path(directory)

    if not directory.exists():
        return []

    return sorted(
        [
            p for p in directory.rglob("*")
            if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS
        ]
    )


# --------------------------------------------------
# CHARGEMENT DOSSIER
# --------------------------------------------------

def load_directory(
    directory: str | Path,
    enable_ocr: bool = False,
    ocr_lang: str = "fra",
) -> str:

    files = list(iter_files(directory))

    if not files:
        return ""

    chunks = []

    for file_path in files:
        try:
            content = load_single_file(
                file_path=file_path,
                enable_ocr=enable_ocr,
                ocr_lang=ocr_lang,
            )

            chunks.append(
                f"\n==================================================\n"
                f"FICHIER : {file_path}\n"
                f"==================================================\n"
                f"{content}\n"
            )

        except Exception as exc:
            chunks.append(
                f"\n==================================================\n"
                f"FICHIER : {file_path}\n"
                f"==================================================\n"
                f"[ERREUR : {exc}]\n"
            )

    return "\n".join(chunks).strip()
